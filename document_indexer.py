import os
import sys
import glob
import json
import re
from typing import List, Dict, Any

# Ép mã hóa UTF-8 cho console Windows
try:
    sys.stdout.reconfigure(encoding='utf-8')
    sys.stderr.reconfigure(encoding='utf-8')
except Exception:
    pass

import pypdf
import docx
import pptx

DATA_DIR = os.path.join(os.path.dirname(__file__), "DuLieu")
INDEX_CACHE_FILE = os.path.join(os.path.dirname(__file__), ".dulieu_index.json")

def clean_text(text: str) -> str:
    if not text:
        return ""
    text = re.sub(r'[ \t]+', ' ', text)
    text = re.sub(r'\n\s*\n+', '\n\n', text)
    return text.strip()

def extract_text_from_pdf(filepath: str, max_pages: int = 150) -> str:
    text = ""
    try:
        reader = pypdf.PdfReader(filepath)
        total_pages = len(reader.pages)
        pages_to_read = min(total_pages, max_pages)
        for i in range(pages_to_read):
            try:
                page = reader.pages[i]
                page_text = page.extract_text()
                if page_text:
                    text += f"\n[Trang {i+1}]\n" + page_text
            except Exception:
                continue
    except Exception as e:
        print(f"  [!] Lỗi đọc PDF '{os.path.basename(filepath)}': {e}")
    return text

def extract_text_from_docx(filepath: str) -> str:
    text = ""
    try:
        doc = docx.Document(filepath)
        for p in doc.paragraphs:
            if p.text.strip():
                text += p.text.strip() + "\n"
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(c.text.strip() for c in row.cells if c.text.strip())
                if row_text:
                    text += row_text + "\n"
    except Exception as e:
        print(f"  [!] Lỗi đọc DOCX '{os.path.basename(filepath)}': {e}")
    return text

def extract_text_from_pptx(filepath: str) -> str:
    text = ""
    try:
        prs = pptx.Presentation(filepath)
        for slide_idx, slide in enumerate(prs.slides):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text.strip():
                            slide_texts.append(paragraph.text.strip())
            if slide_texts:
                text += f"\n[Slide {slide_idx+1}]\n" + "\n".join(slide_texts) + "\n"
    except Exception as e:
        print(f"  [!] Lỗi đọc PPTX '{os.path.basename(filepath)}': {e}")
    return text

def extract_text_from_txt(filepath: str) -> str:
    encodings = ['utf-8', 'utf-16', 'latin-1', 'cp1252']
    for enc in encodings:
        try:
            with open(filepath, "r", encoding=enc) as f:
                return f.read()
        except Exception:
            continue
    return ""

def scan_and_chunk_documents(data_dir: str = DATA_DIR, chunk_size: int = 1500, overlap: int = 200) -> List[Dict[str, Any]]:
    chunks = []
    supported_exts = {'.pdf', '.docx', '.pptx', '.txt', '.md'}
    
    print(f"[*] Đang quét kho tài liệu: {data_dir}...")
    file_count = 0
    
    for root, _, files in os.walk(data_dir):
        for file in files:
            ext = os.path.splitext(file)[1].lower()
            if ext not in supported_exts:
                continue
            if file.startswith("~$"):  # Bỏ qua file tạm của Office
                continue
                
            filepath = os.path.join(root, file)
            rel_path = os.path.relpath(filepath, data_dir)
            file_count += 1
            print(f"  -> Đang xử lý ({file_count}): {rel_path}...")
            
            raw_text = ""
            if ext == '.pdf':
                raw_text = extract_text_from_pdf(filepath)
            elif ext == '.docx':
                raw_text = extract_text_from_docx(filepath)
            elif ext == '.pptx':
                raw_text = extract_text_from_pptx(filepath)
            elif ext in {'.txt', '.md'}:
                raw_text = extract_text_from_txt(filepath)
                
            raw_text = clean_text(raw_text)
            if not raw_text:
                continue
                
            for i in range(0, len(raw_text), chunk_size - overlap):
                chunk_content = raw_text[i:i + chunk_size]
                if len(chunk_content.strip()) < 50:
                    continue
                chunks.append({
                    "source": rel_path,
                    "filename": file,
                    "chunk_id": len(chunks) + 1,
                    "content": chunk_content
                })
                
    print(f"[+] Hoàn tất quét {file_count} tài liệu! Đã tạo {len(chunks)} phân đoạn kiến thức.")
    return chunks

def build_or_load_index(force_refresh: bool = False) -> List[Dict[str, Any]]:
    if not force_refresh and os.path.exists(INDEX_CACHE_FILE):
        try:
            with open(INDEX_CACHE_FILE, "r", encoding="utf-8") as f:
                data = json.load(f)
                print(f"[+] Đã nạp {len(data)} đoạn kiến thức từ bộ nhớ đệm (Cache).")
                return data
        except Exception:
            pass
            
    chunks = scan_and_chunk_documents()
    try:
        with open(INDEX_CACHE_FILE, "w", encoding="utf-8") as f:
            json.dump(chunks, f, ensure_ascii=False, indent=2)
        print(f"[+] Đã lưu bộ nhớ đệm kiến thức vào: {INDEX_CACHE_FILE}")
    except Exception as e:
        print(f"[!] Không thể lưu cache: {e}")
        
    return chunks

def search_relevant_knowledge(query: str, chunks: List[Dict[str, Any]], top_k: int = 5) -> str:
    if not chunks or not query:
        return ""
        
    raw_terms = re.findall(r'\w+', query.lower())
    stop_words = {'là', 'gì', 'như', 'thế', 'nào', 'của', 'và', 'các', 'cho', 'trong', 'được', 'với', 'có', 'the', 'a', 'an', 'is', 'what', 'how', 'cho', 'mình', 'bạn', 'hãy'}
    terms = [t for t in raw_terms if len(t) > 1 and t not in stop_words]
    
    if not terms:
        terms = raw_terms
        
    scored_chunks = []
    for c in chunks:
        content_lower = c["content"].lower()
        source_lower = c["source"].lower()
        score = 0
        
        for t in terms:
            if t in source_lower:
                score += 5
            count = content_lower.count(t)
            score += min(count, 10)
            
        if score > 0:
            scored_chunks.append((score, c))
            
    scored_chunks.sort(key=lambda x: x[0], reverse=True)
    top_results = scored_chunks[:top_k]
    
    if not top_results:
        return ""
        
    combined = ""
    for score, item in top_results:
        combined += f"\n--- [Nguồn: {item['source']}] ---\n{item['content']}\n"
        
    return combined

if __name__ == "__main__":
    index = build_or_load_index(force_refresh=True)
    print(f"Tổng số chunks: {len(index)}")
